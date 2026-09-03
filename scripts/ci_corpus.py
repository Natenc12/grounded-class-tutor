"""Generate a small, deterministic PDF + PPTX corpus so the Slice 2 exit gate can run where the
real one cannot.

The dogfood corpus under `data/dogfood/` is Nate's own course materials and is gitignored, so a
CI runner has no `.pdf`/`.pptx` to enqueue and `ingest_smoke.py` stops at setup. That gate checks
the WRITE PATH — queued → processing → ready, a corrupt file landing `failed(unparseable)`, a
lease reclaim — and none of that depends on what the documents say. So CI gets a generated
corpus: a few multi-page PDFs and slide decks of plain prose, long enough to chunk, built with
the same libraries the ingest tests use for their fixtures (`reportlab`, `python-pptx`).

The Slice 1 gate (`ask_smoke.py`) is different: its questions are anchored to the real corpus
by file and page, and a synthetic one cannot answer them. It stays local-only.

    uv run python scripts/ci_corpus.py <out-dir> [--pdfs 3] [--pptx 2] [--pages 4]

Exit 0 and prints the files written. Deterministic: the same arguments produce the same bytes
of text (the PDF/PPTX containers carry timestamps, so the files themselves are not bit-identical).
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

# Plain prose on a neutral subject. Each entry is one paragraph; pages and slides draw from
# them in order, so a page carries real sentences rather than a repeated token.
PARAGRAPHS = [
    "Water moves through the environment in a continuous cycle driven by the sun. Heat "
    "evaporates water from oceans, lakes and soil, and plants release vapour through their "
    "leaves in a process called transpiration. The vapour rises, cools, and condenses into "
    "droplets that gather as clouds.",
    "When droplets grow heavy enough they fall as precipitation: rain, snow, sleet or hail "
    "depending on the temperature of the air they fall through. Some of that water runs off "
    "the surface into streams and rivers, some soaks into the ground, and some is caught by "
    "vegetation and evaporates again before it ever reaches the soil.",
    "Groundwater is water held in the pores of rock and sediment below the surface. It moves "
    "slowly, sometimes taking centuries to travel a few kilometres, and it feeds springs, "
    "wells and the base flow of rivers during dry seasons. An aquifer is a body of rock "
    "porous enough to store and transmit useful quantities of it.",
    "A watershed is the area of land that drains to a single outlet. Ridges divide one "
    "watershed from the next, and everything that happens on the land inside one, from "
    "farming to paving, changes what arrives at the outlet and when. Engineers use the "
    "boundary to model floods and to plan where a reservoir will fill.",
    "Residence time is how long, on average, a water molecule stays in one part of the "
    "cycle. In the atmosphere it is about nine days. In a large lake it can be years, in "
    "the deep ocean thousands of years, and in some glaciers and aquifers longer still. The "
    "figure matters because it sets how quickly a pollutant clears or a drought ends.",
    "Human use draws mostly on the fast parts of the cycle: rivers, shallow groundwater and "
    "reservoirs. Where withdrawals exceed recharge, water tables fall, wells go dry and land "
    "can subside. Managing that balance is a question of measurement as much as policy, "
    "since a basin cannot be budgeted without knowing what flows in and out of it.",
    "Climate shifts the cycle rather than adding or removing water from it. Warmer air "
    "holds more vapour, so storms carry more and dry spells between them last longer. The "
    "same total rainfall arriving in fewer, heavier events runs off faster and recharges "
    "less, which is why flood and drought risk can rise together in the same region.",
    "Measuring the cycle relies on rain gauges, stream gauges, soil probes and satellites "
    "that weigh the change in gravity as water moves. No single instrument sees the whole "
    "system, so hydrologists close the water balance by combining them and treating the "
    "residual as the sum of what was not measured.",
]


def page_text(page: int, per_page: int = 2) -> list[str]:
    start = (page * per_page) % len(PARAGRAPHS)
    return [PARAGRAPHS[(start + i) % len(PARAGRAPHS)] for i in range(per_page)]


def write_pdf(path: Path, pages: int, title: str) -> None:
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(path), pagesize=LETTER, title=title)
    flow = []
    for p in range(pages):
        flow.append(Paragraph(f"{title} — page {p + 1}", styles["Heading2"]))
        for para in page_text(p):
            flow.append(Paragraph(para, styles["BodyText"]))
            flow.append(Spacer(1, 10))
        if p < pages - 1:
            flow.append(PageBreak())
    doc.build(flow)


def write_pptx(path: Path, slides: int, title: str) -> None:
    deck = Presentation()
    layout = deck.slide_layouts[6]  # blank
    for s in range(slides):
        slide = deck.slides.add_slide(layout)
        box = slide.shapes.add_textbox(0, 0, deck.slide_width, deck.slide_height)
        box.text_frame.text = f"{title} — slide {s + 1}\n\n" + "\n\n".join(page_text(s, 1))
    deck.save(str(path))


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("out_dir")
    ap.add_argument("--pdfs", type=int, default=3)
    ap.add_argument("--pptx", type=int, default=2)
    ap.add_argument("--pages", type=int, default=4, help="pages per PDF and slides per deck")
    args = ap.parse_args()

    out = Path(args.out_dir)
    out.mkdir(parents=True, exist_ok=True)
    written = []
    for i in range(args.pdfs):
        p = out / f"hydrology-lecture-{i + 1:02d}.pdf"
        write_pdf(p, args.pages, f"Hydrology lecture {i + 1}")
        written.append(p)
    for i in range(args.pptx):
        p = out / f"hydrology-slides-{i + 1:02d}.pptx"
        write_pptx(p, args.pages, f"Hydrology slides {i + 1}")
        written.append(p)
    for p in written:
        print(f"{p}  {p.stat().st_size} bytes")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
