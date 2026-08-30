"""Parse - turn a staged file into text + structural units, birthing the source
metadata that the citation spine trusts from here on: citation-spine honor-point ①
(design/decisions/0019-chunking-contract-never-span.md, F2).

Pure function: no DB, no job/status/lease state (the PM-4 seam - Slice 2 wraps this
in worker/queue machinery without rewriting it). Terminal failures (unparseable /
password-protected / zero-text) are signalled by raising `ParseError` with a `reason`
drawn from the same terminal taxonomy as `files.failed_reason`
(design/decisions/0020-ingestion-failure-idempotency-model.md); the caller decides
retry policy, this module never does.

Tooling (pypdf / python-pptx) is a spike choice (design/components/ingestion-worker.md
§"spec/spike line") - swappable without changing the `ParsedUnit` contract.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

import pypdf.errors
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

# The closed terminal-reason taxonomy, mirroring `files.failed_reason`'s CHECK
# (migrations/0001_init.sql, widened by 0003). NOT every value is born here: `too_long` is raised
# by the pipeline's input ceiling, downstream of parsing (ADR 0020, terminal set extended per
# ADR 0029). Every entry is a legal `ParseError` reason; `parse_file` is simply not the only
# place one is raised.
TERMINAL_REASONS = ("unparseable", "protected", "unsupported", "empty", "too_long")

# MS-CFB (OLE2) container signature. Password-protected OOXML files (.pptx/.docx/.xlsx)
# are wrapped in this container instead of being a plain zip, so python-pptx's zip-based
# reader never gets far enough to raise anything more specific - check for it up front.
_OLE_SIGNATURE = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


class ParseError(Exception):
    """A terminal (no-retry) parse failure.

    `reason` matches `files.failed_reason`'s terminal taxonomy (ADR 0020, terminal set
    extended per ADR 0029) so a future caller can pass it straight through with no
    translation. That pass-through is not aspirational any more: `worker.py`'s
    `except ParseError` writes `exc.reason` into the column untranslated.
    """

    def __init__(self, reason: str, message: str) -> None:
        assert reason in TERMINAL_REASONS, f"unknown terminal reason: {reason!r}"
        super().__init__(message)
        self.reason = reason


@dataclass(frozen=True)
class ParsedUnit:
    """One page (PDF) or slide (PPTX) of extracted text.

    Carries the `(file, page_or_slide)` provenance the whole citation spine is born
    from - honor-point ①. `page_or_slide` is 1-indexed to match how a human would cite
    the source ("slide 3"), and is never a range (ADR 0019 never-span).
    """

    text: str
    file: str
    page_or_slide: int


def parse_file(path: str | Path) -> list[ParsedUnit]:
    """Parse a staged PDF or PPTX file into provenance-stamped text units.

    Raises `ParseError` (terminal, no retry) for an unsupported extension, an
    unparseable/corrupt file, a password-protected file, or a file with zero
    extractable text across every page/slide.
    """
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        units = _parse_pdf(path)
    elif suffix == ".pptx":
        units = _parse_pptx(path)
    else:
        raise ParseError("unsupported", f"unsupported file type: {suffix or '(none)'}")

    units = _strip_nul(units)

    if not units:
        raise ParseError("empty", f"no extractable text in {path.name}")

    return units


def _strip_nul(units: list[ParsedUnit]) -> list[ParsedUnit]:
    """Drop NUL (0x00) bytes from every unit's text, and drop units with nothing else left.

    NUL is never legitimate course-material text - it is extraction junk (the real carrier: OCR
    debris on `Livingston Cosmogony.pdf` p.11 of the dogfood corpus). It matters because Postgres
    `text` columns reject it outright: `index_file`'s insert dies with `psycopg.DataError:
    PostgreSQL text fields cannot contain NUL`, AFTER the whole file was parsed, chunked, and
    embedded - money spent, transaction rolled back, file unindexable. Scrubbing belongs HERE, at
    the chokepoint where provenance is born, so every downstream consumer (chunker, embedder,
    index, retrieval, the model's context) sees the same clean text; scrubbing at the DB boundary
    instead would stamp the corpus with text that differs from what was embedded.

    Applied in `parse_file` AFTER format dispatch - one rule for both formats - though only PDFs
    can actually carry it: PPTX bodies are XML 1.0, which cannot represent NUL at all. A unit that
    is nothing but NULs (with whitespace) is dropped exactly like an empty page, and a whole file
    of them lands on the existing `empty` terminal (ADR 0020's taxonomy, no new failure kind).
    """
    cleaned: list[ParsedUnit] = []
    for unit in units:
        text = unit.text.replace("\x00", "")
        if text.strip():
            cleaned.append(
                ParsedUnit(text=text, file=unit.file, page_or_slide=unit.page_or_slide)
                if text != unit.text
                else unit
            )
    return cleaned


def _parse_pdf(path: Path) -> list[ParsedUnit]:
    try:
        reader = PdfReader(path)
    except pypdf.errors.PyPdfError as exc:
        raise ParseError("unparseable", f"could not read PDF {path.name}: {exc}") from exc

    if reader.is_encrypted:
        raise ParseError("protected", f"{path.name} is password-protected")

    units: list[ParsedUnit] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:
            raise ParseError(
                "unparseable", f"could not extract text from {path.name} page {i}: {exc}"
            ) from exc
        if text.strip():
            units.append(ParsedUnit(text=text, file=path.name, page_or_slide=i))
    return units


# Separator between a slide's body text and its speaker notes within the same
# ParsedUnit. Deliberately NOT bracketed. The reason is a nudge, not a collision: the
# Grounder's validator only parses [S#] tokens out of the MODEL'S ANSWER and range-checks
# them (design/decisions/0015-grounder-citation-contract-validation.md §②/§③), so a
# "[Speaker notes]" token in context text would never enter that ladder at all. What it
# WOULD do is put bracketed tokens in front of the model in the very context where [S#] is
# the citation vocabulary — and bracket-shaped context invites bracket-shaped output. Cheap
# to avoid, so avoid it.
_NOTES_MARKER = "Speaker notes:"


def _iter_shapes(shapes):
    """Flatten grouped shapes so text nested inside a PowerPoint group (a common
    authoring pattern) isn't invisible to a top-level `shape.has_text_frame` scan."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _slide_notes(slide) -> str:
    """Return a slide's speaker-notes text, or "" if it has none.

    Guards on `has_notes_slide` FIRST: reading `slide.notes_slide` *creates* the notes
    part as a side effect (verified on python-pptx 1.0.2), so touching it before the
    check both mutates the in-memory deck and makes a "slide has no notes" test pass for
    the wrong reason. `notes_text_frame` can also be None on decks from other authoring
    tools, so it is guarded separately rather than assumed present.
    """
    if not slide.has_notes_slide:
        return ""
    notes_text_frame = slide.notes_slide.notes_text_frame
    if notes_text_frame is None:
        return ""
    return notes_text_frame.text


def _parse_pptx(path: Path) -> list[ParsedUnit]:
    with open(path, "rb") as f:
        header = f.read(len(_OLE_SIGNATURE))
    if header == _OLE_SIGNATURE:
        raise ParseError("protected", f"{path.name} is password-protected")

    try:
        deck = Presentation(str(path))
    except Exception as exc:
        raise ParseError("unparseable", f"could not read PPTX {path.name}: {exc}") from exc

    units: list[ParsedUnit] = []
    for i, slide in enumerate(deck.slides, start=1):
        try:
            lines = []
            for shape in _iter_shapes(slide.shapes):
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line:
                        lines.append(line)
            text = "\n".join(lines)
        except Exception as exc:
            # Body-text failure stays TERMINAL - unchanged by #12.
            raise ParseError(
                "unparseable", f"could not extract text from {path.name} slide {i}: {exc}"
            ) from exc

        try:
            notes = _slide_notes(slide)
        except Exception:
            # DEGRADE, don't fail (D2): a corrupt notes part on an otherwise readable
            # deck must not turn into a terminal `unparseable` the student sees as a
            # refusal. Swallowed deliberately and silently - no logging in V1; Slice 2
            # owns observability (D3).
            notes = ""

        if notes.strip():
            # Body, blank line, marker, notes (D1). Join only the non-empty parts so a
            # notes-only slide doesn't carry a pointless leading blank block.
            parts = [p for p in (text, f"{_NOTES_MARKER}\n{notes}") if p.strip()]
            text = "\n\n".join(parts)

        # Notes ride inside this slide's own unit - one unit per slide, so never-span
        # (ADR 0019) holds by construction. The emptiness check tests the COMBINED text,
        # so a notes-only slide still emits a unit (`empty` is a WHOLE-FILE terminal,
        # ADR 0020).
        if text.strip():
            units.append(ParsedUnit(text=text, file=path.name, page_or_slide=i))
    return units
