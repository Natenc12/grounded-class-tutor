"""Fixture builders for ingest tests.

Real files are generated on the fly (not committed as binary blobs) so `pypdf` /
`python-pptx` / `reportlab` stay the single source of truth for what a "real" PDF or
PPTX looks like.

Also home to `fake_embedder` (issue #4) — a deterministic `Embeddings`-shaped stub so `compose`
can be tested with no network / no OpenAI (the same `client=`-injection discipline as the adapter).

The `db` fixture moved to the root `tests/conftest.py` when the retriever suite (#5) needed it too;
it is still available here by normal conftest inheritance.
"""

from __future__ import annotations

import hashlib
import io
from collections.abc import Sequence
from pathlib import Path

import pytest
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from reportlab.pdfgen import canvas

from gct.config import EMBEDDING_DIM
from gct.ingest.parse import _OLE_SIGNATURE


def make_pdf(path: Path, page_texts: list[str | None], encrypt_password: str | None = None) -> Path:
    """Write a PDF at `path` with one page per entry in `page_texts`.

    A `None` entry produces a genuinely blank page (no text drawn at all), so the
    zero-text test case is a real absence of content, not just whitespace.
    """
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=(612, 792))
    for text in page_texts:
        if text is not None:
            c.drawString(72, 700, text)
        c.showPage()
    c.save()
    buf.seek(0)

    if encrypt_password is None:
        path.write_bytes(buf.getvalue())
        return path

    reader = PdfReader(buf)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(user_password=encrypt_password)
    with open(path, "wb") as f:
        writer.write(f)
    return path


def make_pptx(
    path: Path,
    slide_texts: list[str | None],
    slide_notes: list[str | None] | None = None,
) -> Path:
    """Write a PPTX at `path` with one slide per entry in `slide_texts`.

    A `None` entry produces a slide with no text-bearing shape at all.

    `slide_notes` (issue #12) optionally attaches speaker notes, positionally matched to
    `slide_texts`; a `None` entry leaves that slide with **no notes part at all** - the
    distinction matters, because `slide.notes_slide` creates the part on access, so a slide
    that was merely given empty notes is not the same fixture as one that never had them.
    Defaults to `None` (no slide gets notes), keeping every pre-#12 caller unchanged.
    """
    if slide_notes is not None and len(slide_notes) != len(slide_texts):
        raise ValueError(
            f"slide_notes has {len(slide_notes)} entries but slide_texts has {len(slide_texts)}"
        )

    deck = Presentation()
    blank_layout = deck.slide_layouts[6]  # "Blank" in the default template
    for i, text in enumerate(slide_texts):
        slide = deck.slides.add_slide(blank_layout)
        if text is not None:
            box = slide.shapes.add_textbox(0, 0, deck.slide_width, deck.slide_height)
            box.text_frame.text = text
        if slide_notes is not None and slide_notes[i] is not None:
            slide.notes_slide.notes_text_frame.text = slide_notes[i]
    deck.save(str(path))
    return path


def make_ole_stub(path: Path) -> Path:
    """A file stamped with the MS-CFB signature real password-protected Office
    files use, standing in for one without needing a real encryption library."""
    path.write_bytes(_OLE_SIGNATURE + b"\x00" * 32)
    return path


@pytest.fixture
def pdf_factory(tmp_path):
    def _make(name: str, page_texts: list[str | None], encrypt_password: str | None = None) -> Path:
        return make_pdf(tmp_path / name, page_texts, encrypt_password)

    return _make


@pytest.fixture
def pptx_factory(tmp_path):
    def _make(
        name: str,
        slide_texts: list[str | None],
        slide_notes: list[str | None] | None = None,
    ) -> Path:
        return make_pptx(tmp_path / name, slide_texts, slide_notes)

    return _make


@pytest.fixture
def ole_stub_factory(tmp_path):
    def _make(name: str) -> Path:
        return make_ole_stub(tmp_path / name)

    return _make


# --- Ingest-pipeline fixtures (issue #4) ------------------------------------------------------


class FakeEmbeddings:
    """Deterministic `Embeddings`-shaped stub — no network, no OpenAI.

    Satisfies the `providers.base.Embeddings` protocol (`model_id`, `dim`, `embed`). `dim`
    defaults to `EMBEDDING_DIM` (1536) so vectors it returns satisfy the `chunks.embedding
    vector(1536)` column in the real-DB index tests. Each text maps to a distinct vector (first
    component = a sha256 digest of the text, taken as a 48-bit big-endian int) so a misalignment
    between input order and stored vectors would show up. The digest is stable across processes
    and across time — no `hash()`, no modulo — so the same text always maps to the same vector,
    everywhere, forever (issue #23).

    NOT usable for retrieval ranking: every vector it returns is a positive scalar multiple of
    the same direction, and cosine ignores magnitude — so the cosine distance between any two of
    them is 0 and rank order would be arbitrary. Fine here (ingest never ranks); the retriever
    suite has its own ranking-capable stub in `tests/gct/retriever/conftest.py` (issue #5).
    """

    def __init__(self, model_id: str = "fake-embed-3", dim: int = EMBEDDING_DIM) -> None:
        self._model_id = model_id
        self._dim = dim

    @property
    def model_id(self) -> str:
        return self._model_id

    @property
    def dim(self) -> int:
        return self._dim

    def embed(self, texts: Sequence[str]) -> list[list[float]]:
        vectors: list[list[float]] = []
        for text in texts:
            h = hashlib.sha256(text.encode()).digest()
            # 48 bits, fits float64's mantissa; no modulo
            seed = float(int.from_bytes(h[:6], "big"))
            vectors.append([seed] + [0.0] * (self._dim - 1))
        return vectors


@pytest.fixture
def fake_embedder() -> FakeEmbeddings:
    return FakeEmbeddings()


# --- Ceiling-hardening stubs (issue #43) ------------------------------------------------------
# Both live here rather than beside one test file because two suites in this package need them:
# `test_pipeline.py` (the guard at the `compose`/`ingest_file` seam) and
# `test_ceiling_through_worker.py` (the same guard reached through `enqueue` -> `process_one`).


class CountingEmbeddings:
    """A pass-through wrapper that RECORDS every batch `embed` was handed.

    The ingest suite already has `_ExplodingEmbedder`, which proves a guard fired before `embed`
    by making the call *loud*. This one proves it by making the call *counted*, and the two are
    not interchangeable. An exploding stub can only answer "was `embed` reached" on paths where
    reaching it crashes, so the assertion is really about which exception escaped; `calls == []`
    asks the question directly - was a paid provider round trip bought - and stays answerable on
    the paths where nothing raises at all. That is the assertion the `files`/`jobs` status columns
    cannot make: a guard moved below `embed` sets every column correctly and still spends money.

    Wraps a delegate rather than reimplementing one, so it is usable on BOTH sides of the ceiling:
    the same object can count a refusal at zero and then embed a file that was accepted.
    """

    def __init__(self, inner: FakeEmbeddings | None = None) -> None:
        self._inner = inner if inner is not None else FakeEmbeddings()
        self.calls: list[list[str]] = []

    @property
    def model_id(self) -> str:
        return self._inner.model_id

    @property
    def dim(self) -> int:
        return self._inner.dim

    def embed(self, texts: Sequence[str]) -> list[list[float]]:
        texts = list(texts)
        # Recorded BEFORE the delegate runs: a call that raises still cost a round trip, so this
        # counts calls MADE, not calls that returned (the same stance `tests/gct/jobs`' own stub
        # takes for its transient-failure counter).
        self.calls.append(texts)
        return self._inner.embed(texts)


class EmbedInputTooLarge(RuntimeError):
    """A provider refusing ONE input as over its per-input cap.

    Deliberately NOT `TransientEmbeddingError`. `openai_provider._TRANSIENT_OPENAI_ERRORS` lists
    only 429 / timeout / 5xx / connection-drop, so an over-cap rejection (a 400) propagates out of
    the adapter untouched, and `worker.process_one` classifies it as neither terminal input nor
    bad luck. Modelling it as transient here would make the pipeline look like it retries, which
    is precisely the behavior these tests are pinning. `RuntimeError` keeps it in the same
    "unclassified, propagates" family the real 400 lands in, without importing `openai`.
    """


class PerInputTokenLimitEmbeddings(CountingEmbeddings):
    """A `CountingEmbeddings` that also refuses any SINGLE input larger than one embedding request
    can carry - the cap the real `text-embedding-3-small` enforces and no stub in this repo had.

    `FakeEmbeddings` embeds anything, of any length, for free. That is the right default for every
    test about ordering or stamping, and it is exactly wrong for the space-free-script blind spot
    (issue #43): the whole question there is what happens when ONE chunk is far too big for one
    request, and a stub with no cap answers "it works fine".

    The budget is stated in CHARACTERS, not tokens, on purpose - the assertion is about what the
    PIPELINE does with a provider refusal, not about the accuracy of any token estimate.
    `text-embedding-3-small` accepts at most 8,191 tokens per input; space-free scripts (CJK,
    Thai, Japanese) run roughly 0.5-1.5 tokens per character, so 8,191 tokens is at MOST 16,382
    characters even at the most generous of those rates. An input longer than that has breached
    the real cap under every tokenization of such a script, which is why this bound can be
    asserted without a paid call. It is a LOWER bound on the real cap, never a model of it: this
    stub is not a place to learn what OpenAI charges.
    """

    MAX_INPUT_CHARS = 8_191 * 2

    # Reachable as an attribute so a test can name the type in `pytest.raises` without importing
    # this module. Two `conftest.py` files in this tree share a basename and neither package has
    # an `__init__.py`, so `from conftest import ...` resolves by sys.path order rather than by
    # location - fine in a scratch script, a trap in the suite. The fixture already hands the test
    # an instance; hanging the type off it keeps one source for both.
    Error = EmbedInputTooLarge

    def embed(self, texts: Sequence[str]) -> list[list[float]]:
        texts = list(texts)
        self.calls.append(texts)
        for text in texts:
            if len(text) > self.MAX_INPUT_CHARS:
                raise EmbedInputTooLarge(
                    f"input of {len(text)} characters exceeds what one embedding request "
                    f"can carry ({self.MAX_INPUT_CHARS} characters)"
                )
        return self._inner.embed(texts)


@pytest.fixture
def counting_embedder() -> CountingEmbeddings:
    return CountingEmbeddings()


@pytest.fixture
def token_limit_embedder() -> PerInputTokenLimitEmbeddings:
    return PerInputTokenLimitEmbeddings()
