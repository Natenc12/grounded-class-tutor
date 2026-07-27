"""Unit tests for `gct.ingest.parse` - pure, no DB / API key needed.

Covers the honor-point ① contract: PDF/PPTX happy paths, per-page/slide unit
boundaries (never-span, ADR 0019), provenance stamping, and each terminal-failure
case (ADR 0020).
"""
from __future__ import annotations

import pytest

from gct.ingest.parse import ParsedUnit, ParseError, _slide_notes, parse_file


class TestPdfHappyPath:
    def test_one_unit_per_page_with_text(self, pdf_factory):
        path = pdf_factory("lecture.pdf", ["Page one content", "Page two content"])

        units = parse_file(path)

        assert [u.page_or_slide for u in units] == [1, 2]
        assert "Page one content" in units[0].text
        assert "Page two content" in units[1].text

    def test_units_are_frozen_parsed_units(self, pdf_factory):
        path = pdf_factory("lecture.pdf", ["Some text"])

        [unit] = parse_file(path)

        assert isinstance(unit, ParsedUnit)
        with pytest.raises(AttributeError):
            unit.text = "mutated"  # frozen dataclass - provenance can't drift post-parse

    def test_provenance_stamped_on_every_unit(self, pdf_factory):
        path = pdf_factory("lecture-3.pdf", ["Alpha", "Beta"])

        units = parse_file(path)

        assert all(u.file == "lecture-3.pdf" for u in units)

    def test_blank_page_skipped_without_renumbering_survivors(self, pdf_factory):
        path = pdf_factory("deck.pdf", ["First", None, "Third"])

        units = parse_file(path)

        assert [u.page_or_slide for u in units] == [1, 3]

    def test_unit_never_spans_pages(self, pdf_factory):
        path = pdf_factory("deck.pdf", ["Alpha content", "Beta content"])

        units = parse_file(path)

        assert "Beta content" not in units[0].text
        assert "Alpha content" not in units[1].text


class TestPptxHappyPath:
    def test_one_unit_per_slide_with_text(self, pptx_factory):
        path = pptx_factory("lecture.pptx", ["Slide one content", "Slide two content"])

        units = parse_file(path)

        assert [u.page_or_slide for u in units] == [1, 2]
        assert "Slide one content" in units[0].text
        assert "Slide two content" in units[1].text

    def test_provenance_stamped_on_every_unit(self, pptx_factory):
        path = pptx_factory("lecture-3.pptx", ["Alpha", "Beta"])

        units = parse_file(path)

        assert all(u.file == "lecture-3.pptx" for u in units)

    def test_blank_slide_skipped_without_renumbering_survivors(self, pptx_factory):
        path = pptx_factory("deck.pptx", ["First", None, "Third"])

        units = parse_file(path)

        assert [u.page_or_slide for u in units] == [1, 3]

    def test_unit_never_spans_slides(self, pptx_factory):
        path = pptx_factory("deck.pptx", ["Alpha content", "Beta content"])

        units = parse_file(path)

        assert "Beta content" not in units[0].text
        assert "Alpha content" not in units[1].text

    def test_text_inside_a_grouped_shape_is_extracted(self, tmp_path):
        from pptx import Presentation as PresentationBuilder

        deck = PresentationBuilder()
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        box1 = slide.shapes.add_textbox(0, 0, 100, 100)
        box1.text_frame.text = "Grouped text one"
        box2 = slide.shapes.add_textbox(0, 200, 100, 100)
        box2.text_frame.text = "Grouped text two"
        slide.shapes.add_group_shape([box1, box2])
        path = tmp_path / "grouped.pptx"
        deck.save(str(path))

        [unit] = parse_file(path)

        assert "Grouped text one" in unit.text
        assert "Grouped text two" in unit.text


class TestTerminalFailures:
    def test_unsupported_extension(self, tmp_path):
        path = tmp_path / "notes.txt"
        path.write_text("plain text file")

        with pytest.raises(ParseError) as exc_info:
            parse_file(path)

        assert exc_info.value.reason == "unsupported"

    def test_unparseable_pdf(self, tmp_path):
        path = tmp_path / "corrupt.pdf"
        path.write_bytes(b"this is not a real pdf file")

        with pytest.raises(ParseError) as exc_info:
            parse_file(path)

        assert exc_info.value.reason == "unparseable"

    def test_unparseable_pptx(self, tmp_path):
        path = tmp_path / "corrupt.pptx"
        path.write_bytes(b"this is not a real pptx file")

        with pytest.raises(ParseError) as exc_info:
            parse_file(path)

        assert exc_info.value.reason == "unparseable"

    def test_password_protected_pdf(self, pdf_factory):
        path = pdf_factory("secret.pdf", ["Confidential content"], encrypt_password="hunter2")

        with pytest.raises(ParseError) as exc_info:
            parse_file(path)

        assert exc_info.value.reason == "protected"

    def test_password_protected_pptx(self, ole_stub_factory):
        path = ole_stub_factory("secret.pptx")

        with pytest.raises(ParseError) as exc_info:
            parse_file(path)

        assert exc_info.value.reason == "protected"

    def test_zero_text_pdf(self, pdf_factory):
        path = pdf_factory("blank.pdf", [None, None])

        with pytest.raises(ParseError) as exc_info:
            parse_file(path)

        assert exc_info.value.reason == "empty"

    def test_zero_text_pptx(self, pptx_factory):
        path = pptx_factory("blank.pptx", [None, None])

        with pytest.raises(ParseError) as exc_info:
            parse_file(path)

        assert exc_info.value.reason == "empty"

    def test_pptx_valid_zip_but_not_a_real_ooxml_package_raises_unparseable(self, tmp_path):
        # A structurally valid zip (unlike test_unparseable_pptx's raw garbage bytes)
        # that isn't a real OOXML package - python-pptx fails deep inside package
        # parsing here, not at the zip-open step, so this exercises a different
        # failure mode than the plain-garbage-bytes case above.
        import zipfile

        path = tmp_path / "not_really_a_deck.pptx"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("[Content_Types].xml", "<Types/>")

        with pytest.raises(ParseError) as exc_info:
            parse_file(path)

        assert exc_info.value.reason == "unparseable"

    def test_malformed_slide_content_raises_unparseable(self, pptx_factory, monkeypatch):
        import pptx.shapes.autoshape

        path = pptx_factory("deck.pptx", ["Some text"])

        def _boom(self):
            raise RuntimeError("simulated corrupt slide part")

        monkeypatch.setattr(pptx.shapes.autoshape.Shape, "has_text_frame", property(_boom))

        with pytest.raises(ParseError) as exc_info:
            parse_file(path)

        assert exc_info.value.reason == "unparseable"


class TestSlideNotes:
    def test_slide_with_no_notes_returns_empty(self, pptx_factory):
        import pptx

        path = pptx_factory("deck.pptx", ["Some text"])
        deck = pptx.Presentation(str(path))
        [slide] = deck.slides

        assert _slide_notes(slide) == ""

    def test_notes_slide_not_created_as_a_side_effect(self, pptx_factory):
        import pptx

        path = pptx_factory("deck.pptx", ["Some text"])
        deck = pptx.Presentation(str(path))
        [slide] = deck.slides

        _slide_notes(slide)

        assert slide.has_notes_slide is False

    def test_slide_with_notes_returns_the_notes_text(self, pptx_factory):
        import pptx

        path = pptx_factory("deck.pptx", ["Some text"], slide_notes=["Walk the third way slowly"])
        deck = pptx.Presentation(str(path))
        [slide] = deck.slides

        assert _slide_notes(slide) == "Walk the third way slowly"

    def test_body_and_notes_merged_into_one_unit(self, pptx_factory):
        path = pptx_factory(
            "deck.pptx",
            ["Aquinas: five ways"],
            slide_notes=["The third way is the one they always miss"],
        )

        [unit] = parse_file(path)

        assert unit.page_or_slide == 1
        body_i = unit.text.index("Aquinas: five ways")
        marker_i = unit.text.index("Speaker notes:")
        notes_i = unit.text.index("The third way is the one they always miss")
        assert body_i < marker_i < notes_i

    def test_notes_marker_present_and_unbracketed(self, pptx_factory):
        path = pptx_factory(
            "deck.pptx",
            ["Aquinas: five ways"],
            slide_notes=["The third way is the one they always miss"],
        )

        [unit] = parse_file(path)

        assert "Speaker notes:" in unit.text
        # Unbracketed on purpose: bracket-shaped context invites bracket-shaped output, in the
        # one place where [S#] is the citation vocabulary (ADR 0015 §②; see parse._NOTES_MARKER).
        assert "[Speaker notes]" not in unit.text

    def test_empty_notes_part_adds_no_marker(self, pptx_factory):
        """A notes part that EXISTS but is empty must not emit a bare, dangling marker.

        Probably the most common shape in a real deck: PowerPoint creates the notes part the
        moment anyone clicks into the notes pane, so "part present, text empty" is the default
        state of a deck someone merely looked at. Distinct from the no-notes case
        (`slide_notes=[None]`), which creates no part at all — here the part is real and
        `_slide_notes` returns "", so only the `notes.strip()` check stands between the student
        and an embedded "Speaker notes:" with nothing after it.

        Whitespace-only notes are the same case and covered alongside it.
        """
        for label, notes in (("empty", ""), ("whitespace", "   \n  ")):
            path = pptx_factory(f"deck-{label}.pptx", ["Body text here"], slide_notes=[notes])

            [unit] = parse_file(path)

            assert unit.text == "Body text here", f"{label} notes leaked into the unit"
            assert "Speaker notes:" not in unit.text

    def test_notes_only_slide_yields_a_unit(self, pptx_factory):
        path = pptx_factory(
            "deck.pptx",
            [None],
            slide_notes=["Notes are the only content on this slide"],
        )

        [unit] = parse_file(path)

        assert unit.page_or_slide == 1
        assert "Notes are the only content on this slide" in unit.text

    def test_file_empty_only_when_no_slide_has_body_or_notes(self, pptx_factory):
        all_empty_path = pptx_factory("blank.pptx", [None, None], slide_notes=[None, None])

        with pytest.raises(ParseError) as exc_info:
            parse_file(all_empty_path)
        assert exc_info.value.reason == "empty"

        notes_only_path = pptx_factory(
            "notes-only.pptx", [None], slide_notes=["The only content is in the notes"]
        )

        units = parse_file(notes_only_path)  # must not raise

        assert len(units) == 1

    def test_corrupt_notes_part_keeps_body_text(self, pptx_factory, monkeypatch):
        import pptx.slide

        path = pptx_factory("deck.pptx", ["Some body text"])

        def _boom(self):
            raise RuntimeError("simulated corrupt notes part")

        monkeypatch.setattr(pptx.slide.Slide, "has_notes_slide", property(_boom))

        [unit] = parse_file(path)  # degrades (D2), does not raise

        assert "Some body text" in unit.text
        assert "Speaker notes:" not in unit.text

    def test_notes_do_not_leak_across_slides(self, pptx_factory):
        path = pptx_factory(
            "deck.pptx",
            ["Slide one body", "Slide two body"],
            slide_notes=["Slide one notes", "Slide two notes"],
        )

        units = parse_file(path)

        assert "Slide one notes" in units[0].text
        assert "Slide one notes" not in units[1].text
        assert "Slide two notes" in units[1].text
        assert "Slide two notes" not in units[0].text


class TestNulSanitization:
    """NUL (0x00) never leaves `parse_file` - the write path's one Postgres-fatal character.

    The real carrier is OCR junk on `Livingston Cosmogony.pdf` p.11 (dogfood corpus): one NUL
    that survived parse and chunking, was embedded (money spent), and then killed `index_file`'s
    insert with `psycopg.DataError: PostgreSQL text fields cannot contain NUL` - all-or-nothing
    rollback, file unindexable. These tests pin the fix at its chokepoint.

    The carrier cannot be synthesized honestly: reportlab substitutes NUL before it reaches the
    PDF bytes (it round-trips as a glyph-box, verified), and PPTX bodies are XML 1.0, which
    cannot represent NUL at all. So the format parser is PATCHED to emit the poisoned units -
    which also pins the sanitizer's LOCATION: it must live in `parse_file`, after format
    dispatch, where every format's output passes through it.
    """

    def _poisoned(self, texts, file="scan.pdf"):
        return [
            ParsedUnit(text=text, file=file, page_or_slide=i)
            for i, text in enumerate(texts, start=1)
        ]

    def test_nul_is_stripped_from_unit_text(self, monkeypatch, tmp_path):
        import gct.ingest.parse as parse_mod

        units = self._poisoned(["clean page", "ocr\x00junk\x00here"])
        monkeypatch.setattr(parse_mod, "_parse_pdf", lambda path: units)
        (tmp_path / "scan.pdf").write_bytes(b"%PDF-")  # dispatch is by suffix; content unused

        out = parse_file(tmp_path / "scan.pdf")

        assert [u.text for u in out] == ["clean page", "ocrjunkhere"]
        # Provenance is untouched by the scrub - same file, same page numbers.
        assert [(u.file, u.page_or_slide) for u in out] == [("scan.pdf", 1), ("scan.pdf", 2)]

    def test_nul_only_unit_is_dropped_like_an_empty_page(self, monkeypatch, tmp_path):
        import gct.ingest.parse as parse_mod

        units = self._poisoned(["\x00\x00 \x00", "real content"])
        monkeypatch.setattr(parse_mod, "_parse_pdf", lambda path: units)
        (tmp_path / "scan.pdf").write_bytes(b"%PDF-")

        out = parse_file(tmp_path / "scan.pdf")

        # The all-NUL page vanishes; the survivor keeps ITS OWN page number (no renumbering),
        # exactly like the existing blank-page behavior.
        assert [(u.text, u.page_or_slide) for u in out] == [("real content", 2)]

    def test_whole_file_of_nul_lands_on_the_empty_terminal(self, monkeypatch, tmp_path):
        import gct.ingest.parse as parse_mod

        units = self._poisoned(["\x00", "\x00 \x00"])
        monkeypatch.setattr(parse_mod, "_parse_pdf", lambda path: units)
        (tmp_path / "scan.pdf").write_bytes(b"%PDF-")

        with pytest.raises(ParseError) as excinfo:
            parse_file(tmp_path / "scan.pdf")

        # The existing `empty` terminal, not a new failure kind (ADR 0020's taxonomy is closed).
        assert excinfo.value.reason == "empty"
